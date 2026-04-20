from app.main.util.extractors.common import normalize_lines
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def read_docx(path: str) -> str:
    """Read all text from a .docx file, including table cells."""

    doc = Document(path)
    texts = []

    for p in doc.paragraphs:
        if p.text.strip():
            texts.append(p.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                texts.append(" | ".join(row_text))

    return "\n".join(texts)


def read_csv(path: str) -> str:
    import csv

    with open(path, "r", encoding="utf-8-sig", errors="ignore", newline="") as f:
        sample = f.read(4096)
        f.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            dialect = csv.excel
        rows = ["\t".join(str(c).strip() for c in row) for row in csv.reader(f, dialect)]

    return normalize_lines(rows)


def read_excel(path: str) -> str:

    wb = load_workbook(path, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            out.append("\t".join("" if v is None else str(v).strip() for v in row))
        out.append("")
    return normalize_lines(out)


def read_pptx(path: str) -> str:
    prs = Presentation(path)
    lines = []

    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
            elif shape.has_table:
                lines.append("\n--- TABLE START ---")
                for row in shape.table.rows:
                    lines.append("\t".join(cell.text_frame.text.strip() for cell in row.cells))
                lines.append("--- TABLE END ---\n")
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                pass

        lines.append("")

    return normalize_lines(lines)
